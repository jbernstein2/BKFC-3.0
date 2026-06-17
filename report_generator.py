# report_generator.py

from typing import Dict
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def _add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def _add_kpi_slide(prs: Presentation, title: str, kpis: Dict[str, float]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    lines = [f"{k}: {round(v, 2)}" for k, v in kpis.items() if v is not None]
    body.text = "\n".join(lines)


def generate_match_deck(
    output_path: str,
    match_kpis_brooklyn: Dict[str, float],
    match_kpis_opponent: Dict[str, float],
    brooklyn_season_avg: Dict[str, float],
    opponent_season_avg: Dict[str, float],
    opponent_name: str,
    match_label: str,
):
    """
    Generate a simple PowerPoint deck summarizing match and season KPIs.
    """
    prs = Presentation()

    _add_title_slide(
        prs,
        title=f"Brooklyn FC Match Report",
        subtitle=match_label,
    )

    _add_kpi_slide(prs, "Match KPIs - Brooklyn", match_kpis_brooklyn)
    _add_kpi_slide(prs, f"Match KPIs - {opponent_name}", match_kpis_opponent)

    _add_kpi_slide(prs, "Season Averages - Brooklyn", brooklyn_season_avg)
    _add_kpi_slide(
        prs,
        f"Season Averages - {opponent_name}",
        opponent_season_avg,
    )

    output_path = Path(output_path)
    prs.save(output_path)
    return str(output_path)
